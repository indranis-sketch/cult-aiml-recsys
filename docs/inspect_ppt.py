import collections.abc
from pptx import Presentation
import os

prs = Presentation('Presentation.pptx')

for i, slide in enumerate(prs.slides):
    print(f"Slide {i+1}:")
    print(f"  Layout: {slide.slide_layout.name}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"  - Shape (Text): {shape.text[:50]}...")
        elif shape.has_image:
            print(f"  - Shape (Image)")
        else:
            print(f"  - Shape (Other: {shape.shape_type})")
