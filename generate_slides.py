import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    prs = Presentation()
    
    slides_data = [
        {
            "title": "Personalized Content Discovery Engine",
            "content": "Cult Open Projects 2026 - AI/ML Track\nPrepared by: [Your Name]"
        },
        {
            "title": "The Problem: Navigating the Sea of Content",
            "content": "• The Challenge: Modern streaming platforms host massive content libraries. Users experience decision paralysis.\n• The Goal: Build an intelligent system that learns individual user preferences and surfaces highly relevant content.\n• The Dataset: The Netflix Prize Dataset—containing millions of ratings."
        },
        {
            "title": "Data Exploration & Sparsity",
            "content": "• The 98% Rule: The dataset is highly sparse. Most users have only interacted with a fraction of a percent of the total movie catalog.\n• The Bias: Ratings skew positive (mostly 3, 4, and 5 stars). Users rarely rate movies they actively avoid watching.\n• The Implications: We cannot rely on simple nearest-neighbor models due to lack of overlapping interactions."
        },
        {
            "title": "The Long Tail Phenomenon",
            "content": "• Content Popularity: A tiny percentage of blockbuster films receive millions of ratings, while thousands of indie films receive almost none.\n• User Activity: A small group of power users generates the bulk of the data.\n• The Goal: A robust engine must look past blockbusters to help users discover the long tail of niche content."
        },
        {
            "title": "Methodology & Architecture",
            "content": "• Data Processing: Raw text data parsed into heavily optimized scipy.sparse matrices to handle millions of interactions efficiently.\n• Train/Test Split: Data divided 80/20 to ensure models are evaluated on unseen user behavior.\n• The Algorithm: Transitioned from baseline Memory-based Collaborative Filtering to advanced Model-based Matrix Factorization."
        },
        {
            "title": "The Model: Singular Value Decomposition",
            "content": "• How it Works: SVD breaks down the massive user-item grid into two smaller matrices: User-Features and Item-Features.\n• Latent Factors: It automatically discovers hidden themes (e.g., Action-Heavy, Dark Comedy) without being explicitly programmed.\n• Prediction: By multiplying a user's hidden preferences by a movie's hidden traits, we can predict their exact rating."
        },
        {
            "title": "Evaluation Metrics",
            "content": "1. RMSE (Root Mean Squared Error):\n   • Measures the mathematical accuracy of our predicted ratings against the real test data. Lower is better.\n\n2. MAP@10 (Mean Average Precision at 10):\n   • Measures the ranking quality of our Top 10 recommendations.\n   • We only rewarded the model if the recommended movie actually received a rating of 3.5 stars or higher."
        },
        {
            "title": "Future Roadmap",
            "content": "• Hybrid Systems: Incorporating movie metadata (Genres, Directors) to solve the Cold Start problem for brand-new users.\n• Deep Learning: Exploring Neural Collaborative Filtering (NCF) to capture non-linear, complex behavioral patterns.\n• Explainable AI: Adding transparency features to the UI (e.g., 'Recommended because you previously watched...') to build trust."
        }
    ]
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = slides_data[0]["title"]
    subtitle.text = slides_data[0]["content"]
    
    # Slides 2-8: Bullet points
    for slide_data in slides_data[1:]:
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data["title"]
        content.text = slide_data["content"]
        
        # Adjust font sizes
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), 'docs', 'Presentation.pptx'))
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
